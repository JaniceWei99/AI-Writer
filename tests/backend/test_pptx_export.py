"""Tests for services/pptx_export.py — Markdown to PPTX conversion."""

import io
import pytest
from pptx import Presentation

from services.pptx_export import markdown_to_pptx, _parse_slides, THEMES, VALID_LAYOUTS


# ── Parser ───────────────────────────────────────────────────────────


class TestParseSlides:
    def test_basic_heading_and_bullets(self):
        md = "## Slide One\n- Bullet A\n- Bullet B"
        slides = _parse_slides(md)
        assert len(slides) == 1
        assert slides[0].title == "Slide One"
        assert slides[0].bullets == ["Bullet A", "Bullet B"]

    def test_numbered_heading(self):
        md = "## 1. Introduction\n- Point"
        slides = _parse_slides(md)
        assert slides[0].title == "Introduction"

    def test_multiple_slides_with_separator(self):
        md = "## First\n- A\n---\n## Second\n- B"
        slides = _parse_slides(md)
        assert len(slides) == 2
        assert slides[0].title == "First"
        assert slides[1].title == "Second"

    def test_multiple_slides_without_separator(self):
        md = "## First\n- A\n## Second\n- B"
        slides = _parse_slides(md)
        assert len(slides) == 2

    def test_speaker_notes(self):
        md = "## Title\n- Bullet\n> Speaker notes here"
        slides = _parse_slides(md)
        assert slides[0].notes == "Speaker notes here"

    def test_speaker_notes_prefix_stripped(self):
        md = "## Title\n- Bullet\n> 演讲备注：This is the note"
        slides = _parse_slides(md)
        assert slides[0].notes == "This is the note"

    def test_ordered_list_as_bullets(self):
        md = "## Title\n1. First\n2. Second"
        slides = _parse_slides(md)
        assert slides[0].bullets == ["First", "Second"]

    def test_empty_input(self):
        slides = _parse_slides("")
        assert slides == []

    def test_only_whitespace(self):
        slides = _parse_slides("   \n  \n  ")
        assert slides == []

    def test_content_before_first_heading(self):
        md = "Some intro text\n## Actual Slide\n- Bullet"
        slides = _parse_slides(md)
        assert len(slides) == 2
        assert slides[0].bullets == ["Some intro text"]
        assert slides[1].title == "Actual Slide"

    def test_asterisk_and_plus_bullets(self):
        md = "## Title\n* Star bullet\n+ Plus bullet"
        slides = _parse_slides(md)
        assert slides[0].bullets == ["Star bullet", "Plus bullet"]

    def test_multi_line_notes(self):
        md = "## Title\n- Bullet\n> Line one\n> Line two"
        slides = _parse_slides(md)
        assert "Line one" in slides[0].notes
        assert "Line two" in slides[0].notes

    # ── Layout parsing ──

    def test_layout_tag_parsed(self):
        md = "## 2. Market Data [layout: stats]\n- Revenue: 500M"
        slides = _parse_slides(md)
        assert slides[0].layout == "stats"
        assert slides[0].title == "Market Data"

    def test_layout_tag_case_insensitive(self):
        md = "## Title [Layout: COMPARISON]\n- A | B"
        slides = _parse_slides(md)
        assert slides[0].layout == "comparison"

    def test_layout_tag_all_types(self):
        for layout in VALID_LAYOUTS:
            md = f"## Test [layout: {layout}]\n- Item"
            slides = _parse_slides(md)
            assert slides[0].layout == layout

    def test_unknown_layout_defaults_to_bullets(self):
        md = "## Test [layout: unknown_type]\n- Item"
        slides = _parse_slides(md)
        assert slides[0].layout == "bullets"
        # Unknown tag text stays in title
        assert "unknown_type" in slides[0].title

    def test_no_layout_tag_defaults_to_bullets(self):
        md = "## Plain Title\n- Bullet"
        slides = _parse_slides(md)
        assert slides[0].layout == "bullets"

    def test_layout_stripped_from_title(self):
        md = "## 3. 核心数据 [layout: stats]\n- 市场规模: 500亿"
        slides = _parse_slides(md)
        assert slides[0].title == "核心数据"
        assert "[layout" not in slides[0].title

    def test_mixed_layouts_in_deck(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## Data [layout: stats]\n- Users: 1M\n---\n"
            "## Compare [layout: comparison]\n- A | B\n---\n"
            "## History [layout: timeline]\n- 2020: Start\n---\n"
            "## Vision [layout: quote]\n- Dream big\n---\n"
            "## Features [layout: grid]\n- F1: Desc\n---\n"
            "## Normal [layout: bullets]\n- Point\n---\n"
            "## 谢谢\n- Q&A"
        )
        slides = _parse_slides(md)
        assert len(slides) == 8
        assert slides[1].layout == "stats"
        assert slides[2].layout == "comparison"
        assert slides[3].layout == "timeline"
        assert slides[4].layout == "quote"
        assert slides[5].layout == "grid"
        assert slides[6].layout == "bullets"


# ── PPTX builder ─────────────────────────────────────────────────────


class TestMarkdownToPptx:
    def test_returns_bytesio(self):
        md = "## Slide 1\n- Hello"
        buf = markdown_to_pptx(md, "Test")
        assert isinstance(buf, io.BytesIO)
        assert buf.getbuffer().nbytes > 0

    def test_valid_pptx(self):
        md = "## Cover\n- Subtitle\n---\n## Content\n- Bullet\n---\n## 谢谢\n- Q&A"
        buf = markdown_to_pptx(md, "Test Presentation")
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_single_slide(self):
        md = "## Only One\n- Point"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 1

    def test_fallback_when_no_slides(self):
        md = ""
        buf = markdown_to_pptx(md, "Fallback Title")
        prs = Presentation(buf)
        assert len(prs.slides) == 1

    def test_all_themes(self):
        md = "## Slide\n- Bullet"
        for template_name in THEMES:
            buf = markdown_to_pptx(md, template=template_name)
            prs = Presentation(buf)
            assert len(prs.slides) >= 1

    def test_unknown_template_falls_back(self):
        md = "## Slide\n- Bullet"
        buf = markdown_to_pptx(md, template="nonexistent")
        prs = Presentation(buf)
        assert len(prs.slides) >= 1

    def test_speaker_notes_embedded(self):
        md = "## Slide\n- Bullet\n> My speaker note"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "My speaker note" in notes_text

    def test_slide_dimensions(self):
        md = "## Slide\n- Bullet"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        # 13.33 inches wide, 7.5 inches tall
        assert prs.slide_width > 0
        assert prs.slide_height > 0

    def test_with_images(self):
        """Content slides with images should have a picture shape."""
        md = "## Cover\n- Subtitle\n---\n## Content Slide\n- Bullet\n---\n## 谢谢\n- Q&A"
        # Create a minimal valid PNG (1x1 pixel)
        import struct, zlib
        def _make_png():
            sig = b'\x89PNG\r\n\x1a\n'
            ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = struct.pack('>I', zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff)
            ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + ihdr_crc
            raw = zlib.compress(b'\x00\x00\x00\x00')
            idat_crc = struct.pack('>I', zlib.crc32(b'IDAT' + raw) & 0xffffffff)
            idat = struct.pack('>I', len(raw)) + b'IDAT' + raw + idat_crc
            iend_crc = struct.pack('>I', zlib.crc32(b'IEND') & 0xffffffff)
            iend = struct.pack('>I', 0) + b'IEND' + iend_crc
            return sig + ihdr + idat + iend

        png_bytes = _make_png()
        # Provide image for slide index 1 (the content slide)
        images = {1: png_bytes}
        buf = markdown_to_pptx(md, images=images)
        prs = Presentation(buf)
        assert len(prs.slides) == 3
        # Content slide (index 1) should have a picture
        content_slide = prs.slides[1]
        from pptx.shapes.picture import Picture
        pictures = [s for s in content_slide.shapes if isinstance(s, Picture)]
        assert len(pictures) == 1

    def test_without_images_no_pictures(self):
        """Without images dict, no picture shapes should be present."""
        md = "## Cover\n- Sub\n---\n## Content\n- Bullet"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        from pptx.shapes.picture import Picture
        for slide in prs.slides:
            pictures = [s for s in slide.shapes if isinstance(s, Picture)]
            assert len(pictures) == 0


# ── Layout-specific rendering ────────────────────────────────────────


class TestLayoutRendering:
    """Verify each layout type generates a valid PPTX without errors."""

    def test_stats_layout(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## 核心数据 [layout: stats]\n"
            "- 市场规模: 500亿\n- 年增长率: 35%\n- 用户数: 2.3亿\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_comparison_layout(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## 方案对比 [layout: comparison]\n"
            "- 传统方案 | AI方案\n"
            "- 人工编写 3天 | 自动生成 3分钟\n"
            "- 成本 5万元 | 成本 500元\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_timeline_layout(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## 发展历程 [layout: timeline]\n"
            "- 2020: 项目启动\n- 2022: 用户百万\n- 2024: 全球化\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_quote_layout(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## 核心理念 [layout: quote]\n"
            "- 让每个人都能用AI创作高质量内容\n"
            "- 创始人 张三\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_grid_layout(self):
        md = (
            "## Cover\n- Subtitle\n---\n"
            "## 核心功能 [layout: grid]\n"
            "- 智能写作: 基于AI的一键生成\n"
            "- PPT制作: 自动排版与设计\n"
            "- 多平台适配: 小红书、公众号等\n"
            "- 数据分析: 智能报表生成\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_all_layouts_in_one_deck(self):
        """Full deck with all layout types generates without errors."""
        md = (
            "## AI写作助手介绍\n- 新一代智能写作平台\n---\n"
            "## 目录 [layout: bullets]\n- 产品概述\n- 核心数据\n- 功能对比\n---\n"
            "## 核心指标 [layout: stats]\n- 用户数: 100万\n- 满意度: 98%\n- 日活: 50万\n---\n"
            "## 传统 vs AI [layout: comparison]\n- 传统写作 | AI写作\n- 耗时长 | 秒级生成\n---\n"
            "## 发展历程 [layout: timeline]\n- 2023: 立项\n- 2024: 上线\n- 2025: 百万用户\n---\n"
            "## 愿景 [layout: quote]\n- 让创作更简单\n- CEO\n---\n"
            "## 功能矩阵 [layout: grid]\n- 写作: AI生成\n- PPT: 自动排版\n- 翻译: 多语言\n---\n"
            "## 谢谢\n- Q&A"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 8

    def test_all_layouts_all_themes(self):
        """Each layout renders successfully with every theme."""
        md = (
            "## Cover\n- Sub\n---\n"
            "## Stats [layout: stats]\n- A: 100\n- B: 200\n---\n"
            "## Compare [layout: comparison]\n- X | Y\n- 1 | 2\n---\n"
            "## Time [layout: timeline]\n- 2020: Start\n---\n"
            "## Quote [layout: quote]\n- Big idea\n---\n"
            "## Grid [layout: grid]\n- Card: Desc\n---\n"
            "## 谢谢\n- Bye"
        )
        for theme_name in THEMES:
            buf = markdown_to_pptx(md, template=theme_name)
            prs = Presentation(buf)
            assert len(prs.slides) == 7

    def test_stats_single_item(self):
        """Stats layout with a single item should not crash."""
        md = "## Cover\n- Sub\n---\n## Data [layout: stats]\n- Revenue: 1B\n---\n## 谢谢\n- End"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_comparison_no_pipe(self):
        """Comparison layout with bullets that lack pipe falls back gracefully."""
        md = (
            "## Cover\n- Sub\n---\n"
            "## Compare [layout: comparison]\n"
            "- Item without pipe\n- Another item\n---\n"
            "## 谢谢\n- End"
        )
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_grid_many_items(self):
        """Grid with more than 6 items caps at 9."""
        items = "\n".join(f"- Item{i}: Desc{i}" for i in range(10))
        md = f"## Cover\n- Sub\n---\n## Big Grid [layout: grid]\n{items}\n---\n## 谢谢\n- End"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3

    def test_quote_single_line(self):
        """Quote layout with a single bullet (no attribution)."""
        md = "## Cover\n- Sub\n---\n## Insight [layout: quote]\n- One powerful idea\n---\n## 谢谢\n- End"
        buf = markdown_to_pptx(md)
        prs = Presentation(buf)
        assert len(prs.slides) == 3
