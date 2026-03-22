# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0.0"]
# ///
"""Generate a project presentation PPT for AI 写作助手."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────
WHITE       = RGBColor(255, 255, 255)
BLACK       = RGBColor(30, 30, 30)
DARK_BG     = RGBColor(24, 24, 36)
ACCENT      = RGBColor(79, 140, 255)
ACCENT_LIGHT= RGBColor(130, 177, 255)
LIGHT_GRAY  = RGBColor(180, 180, 195)
MEDIUM_GRAY = RGBColor(120, 120, 140)
CARD_BG     = RGBColor(36, 36, 54)
GREEN       = RGBColor(74, 222, 128)
ORANGE      = RGBColor(251, 191, 36)
PURPLE      = RGBColor(167, 139, 250)
PINK        = RGBColor(244, 114, 182)
CYAN        = RGBColor(34, 211, 238)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def _set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # bullet
        run_b = p.add_run()
        run_b.text = "●  "
        run_b.font.size = Pt(font_size - 2)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = "Microsoft YaHei"
        # text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Microsoft YaHei"
    return txBox


def _add_card(slide, left, top, width, height, title, body_lines,
              accent=ACCENT, bg=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False

    # accent top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    _add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                  title, font_size=18, color=WHITE, bold=True)
    _add_bullet_list(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                     height - Inches(0.9), body_lines, font_size=13, color=LIGHT_GRAY,
                     bullet_color=accent)


def _add_section_title(slide, text, subtitle=""):
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                  text, font_size=36, color=WHITE, bold=True)
    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
                      subtitle, font_size=16, color=MEDIUM_GRAY)
    # accent underline
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(1.2), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _add_arch_box(slide, left, top, w, h, label, sub, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    _add_text_box(slide, left, top + Inches(0.15), w, Inches(0.4),
                  label, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, left, top + Inches(0.55), w, Inches(0.4),
                  sub, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


def _add_arrow(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ══════════════════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)

_add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
              "AI 写作助手", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
              "基于本地大语言模型的智能写作工具", font_size=24, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(1), Inches(4.1), Inches(11), Inches(0.5),
              "FastAPI  ·  React  ·  Ollama  ·  Qwen 3.5", font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
              "2026.03", font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 2: 项目概念
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "项目概念", "What & Why")

_add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
              "AI 写作助手是一个完全本地化的智能写作平台，通过 Ollama 运行开源大语言模型，在保障数据隐私的前提下，"
              "为用户提供文章生成、润色、翻译、摘要等一站式写作服务。",
              font_size=18, color=LIGHT_GRAY)

_add_card(slide, Inches(0.8), Inches(3.0), Inches(3.6), Inches(3.2),
          "核心价值", [
              "数据隐私：完全本地运行，不依赖云端 API",
              "多场景覆盖：8 种写作风格 + 诗词创作",
              "实时体验：SSE 流式输出，逐字显示",
              "开放灵活：可更换任意 Ollama 模型",
          ], accent=ACCENT)

_add_card(slide, Inches(4.8), Inches(3.0), Inches(3.6), Inches(3.2),
          "目标用户", [
              "内容创作者（自媒体、新媒体运营）",
              "学生和研究人员（论文、翻译）",
              "开发者（了解 LLM 应用架构）",
              "注重数据隐私的企业用户",
          ], accent=GREEN)

_add_card(slide, Inches(8.8), Inches(3.0), Inches(3.6), Inches(3.2),
          "功能矩阵", [
              "文章生成 / 文本润色 / 翻译 / 摘要",
              "小红书 / 公众号 / 头条 / 短剧",
              "古典诗词（五言 / 七言）",
              "附件上传解析 / Word 导出",
          ], accent=ORANGE)


# ══════════════════════════════════════════════════════════════
# Slide 3: 系统架构
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "系统架构", "Architecture Overview")

# Three boxes
_add_arch_box(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(1.2),
              "React 前端", "Vite + TypeScript · :5173", ACCENT)
_add_arrow(slide, Inches(4.2), Inches(2.85), Inches(1.2))
_add_text_box(slide, Inches(4.2), Inches(3.3), Inches(1.2), Inches(0.3),
              "HTTP / SSE", font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

_add_arch_box(slide, Inches(5.6), Inches(2.5), Inches(3.2), Inches(1.2),
              "FastAPI 后端", "Python 3.10+ · :8000", GREEN)
_add_arrow(slide, Inches(9.0), Inches(2.85), Inches(1.2))
_add_text_box(slide, Inches(9.0), Inches(3.3), Inches(1.2), Inches(0.3),
              "HTTP Stream", font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

_add_arch_box(slide, Inches(10.4), Inches(2.5), Inches(2.2), Inches(1.2),
              "Ollama", "qwen3.5:9b · :11434", PURPLE)

# Bottom detail cards
_add_card(slide, Inches(0.8), Inches(4.4), Inches(3.6), Inches(2.5),
          "前端层", [
              "React 组件：表单 / 结果 / 历史",
              "Axios + 原生 Fetch (SSE)",
              "Markdown 实时渲染",
              "localStorage 历史记录",
          ], accent=ACCENT)

_add_card(slide, Inches(4.8), Inches(4.4), Inches(3.6), Inches(2.5),
          "后端层", [
              "路由分发 + 请求校验 (Pydantic)",
              "提示词模板引擎 (9 种模板)",
              "文件解析 (PDF / Word / 纯文本)",
              "Word 导出 (Markdown → .docx)",
          ], accent=GREEN)

_add_card(slide, Inches(8.8), Inches(4.4), Inches(3.6), Inches(2.5),
          "模型层", [
              "Ollama REST API 调用",
              "流式 / 非流式两种模式",
              "诗词格律校验 + 自动重试",
              "本地推理，零网络依赖",
          ], accent=PURPLE)


# ══════════════════════════════════════════════════════════════
# Slide 4: 后端技术栈
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "后端技术栈", "Backend: Python + FastAPI")

items = [
    ("FastAPI", "高性能异步 Web 框架，自动生成 OpenAPI 文档，原生支持 async/await 和流式响应", ACCENT),
    ("Pydantic v2", "数据校验与序列化，定义 WritingRequest / WritingResponse / ExportRequest 等模型", GREEN),
    ("httpx", "异步 HTTP 客户端，用于调用 Ollama API，支持流式读取响应体 (aiter_lines)", ORANGE),
    ("python-docx", "Word 文档生成库，将 Markdown 转换为带格式的 .docx 文件（标题/列表/加粗/斜体）", PURPLE),
    ("PyPDF2", "PDF 文本提取，逐页读取并合并文本内容", PINK),
    ("uvicorn", "ASGI 服务器，支持多 worker 生产部署，可搭配 gunicorn 使用", CYAN),
]

for i, (title, desc, color) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(1.7)

    _add_text_box(slide, x, y, Inches(5.8), Inches(0.4),
                  title, font_size=20, color=color, bold=True)
    _add_text_box(slide, x, y + Inches(0.4), Inches(5.8), Inches(0.8),
                  desc, font_size=14, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 5: 前端技术栈
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "前端技术栈", "Frontend: React + TypeScript + Vite")

items = [
    ("React 19", "函数式组件 + Hooks 状态管理，useState / useRef / useCallback / useMemo", ACCENT),
    ("TypeScript", "完整类型定义，TaskType / WritingRequest / HistoryItem 等接口契约", GREEN),
    ("Vite 8", "极速开发服务器 (HMR)，生产构建输出优化的静态文件", ORANGE),
    ("Axios + Fetch", "Axios 处理常规请求，原生 Fetch + ReadableStream 处理 SSE 流式响应", PURPLE),
    ("react-markdown", "Markdown → React 组件的实时渲染，支持标题/列表/代码块/引用等", PINK),
    ("File System Access API", 'showSaveFilePicker 弹出系统"另存为"对话框, 支持用户选择路径和重命名', CYAN),
]

for i, (title, desc, color) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(1.7)

    _add_text_box(slide, x, y, Inches(5.8), Inches(0.4),
                  title, font_size=20, color=color, bold=True)
    _add_text_box(slide, x, y + Inches(0.4), Inches(5.8), Inches(0.8),
                  desc, font_size=14, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 6: 大模型调用流程
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "大模型调用流程", "How We Call the LLM")

steps = [
    ("1", "用户输入", "选择任务类型\n输入文本/主题", ACCENT),
    ("2", "构建 Prompt", "build_prompt()\n选择模板 + 注入参数", GREEN),
    ("3", "调用 Ollama", "POST /api/generate\nstream=true, think=false", ORANGE),
    ("4", "流式返回", "SSE: data: \"token\"\n逐 token 推送", PURPLE),
    ("5", "渲染展示", "Markdown 实时渲染\n保存历史记录", PINK),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.5)
    y = Inches(2.2)

    # Circle number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    _add_text_box(slide, x + Inches(0.7), y + Inches(0.08), Inches(0.6), Inches(0.5),
                  num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, x, y + Inches(0.8), Inches(2.1), Inches(0.4),
                  title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, x, y + Inches(1.2), Inches(2.1), Inches(0.8),
                  desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 4:
        _add_text_box(slide, x + Inches(2.05), y + Inches(0.12), Inches(0.5), Inches(0.5),
                      "→", font_size=24, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: code-like explanation
code_text = (
    "# 后端核心调用链\n"
    "prompt = build_prompt(task_type, content, style, target_lang, attachment_text)\n"
    "async for token in generate_stream(prompt):       # httpx 流式调用 Ollama\n"
    "    yield f'data: {json.dumps(token)}\\n\\n'       # SSE 格式推送给前端"
)
code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8))
code_box.fill.solid()
code_box.fill.fore_color.rgb = RGBColor(20, 20, 30)
code_box.line.color.rgb = RGBColor(60, 60, 80)
code_box.line.width = Pt(1)

_add_text_box(slide, Inches(1.1), Inches(4.95), Inches(11.2), Inches(1.5),
              code_text, font_size=14, color=GREEN, font_name="Consolas")


# ══════════════════════════════════════════════════════════════
# Slide 7: Ollama 与模型架构
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "Ollama 与模型架构", "Ollama Runtime & Qwen Model")

_add_card(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
          "Ollama 运行时", [
              "本地模型推理引擎，类似 Docker 管理容器镜像的方式管理 LLM",
              "提供 REST API (localhost:11434)，兼容 OpenAI 接口格式",
              "自动管理 GPU/CPU 调度、模型加载卸载、显存分配",
              "支持 GGUF 量化格式，降低显存需求 (9B 模型约需 6GB VRAM)",
          ], accent=PURPLE)

_add_card(slide, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.4),
          "Qwen 3.5 (9B) 模型", [
              "阿里通义千问团队开源的大语言模型",
              "9B 参数量，平衡性能与资源消耗",
              "中英双语优化，中文写作质量优秀",
              "支持 32K 上下文长度，适合长文生成",
          ], accent=ORANGE)

_add_card(slide, Inches(0.8), Inches(4.6), Inches(11.6), Inches(2.5),
          "调用方式对比", [
              "非流式 (stream=false): POST /api/generate → 等待完整生成 → 一次返回 {response, eval_count}",
              "流式 (stream=true): POST /api/generate → 逐行返回 JSON → 每行含 {response: token, done: bool}",
              "think=false: 禁用思考链 (Chain-of-Thought) 输出，避免 <think> 标签混入结果",
              "超时设置 120 秒: 长文本 (1000+ tokens) 生成可能需要较长时间",
          ], accent=ACCENT)


# ══════════════════════════════════════════════════════════════
# Slide 8: 大语言模型原理
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "大语言模型原理", "How LLMs Work")

_add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.2),
          "Transformer 架构", [
              "自注意力机制 (Self-Attention)",
              "并行处理序列中所有位置的关系",
              "多头注意力捕获不同语义维度",
              "位置编码保留序列顺序信息",
          ], accent=ACCENT)

_add_card(slide, Inches(4.8), Inches(1.8), Inches(3.6), Inches(2.2),
          "自回归生成", [
              "逐 token 生成: P(x_t | x_1...x_{t-1})",
              "每步从概率分布中采样下一个 token",
              "Temperature 控制随机性 (创造性)",
              "Top-p / Top-k 截断低概率候选",
          ], accent=GREEN)

_add_card(slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(2.2),
          "训练过程", [
              "预训练: 海量文本的下一词预测",
              "SFT: 指令微调，学习遵循指令",
              "RLHF / DPO: 人类偏好对齐",
              "结果: 理解意图 + 高质量输出",
          ], accent=ORANGE)

_add_card(slide, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.6),
          "Prompt Engineering (本项目的应用)", [
              "角色设定: 每个模板定义专家角色 (\"你是一位专业的中文写作助手\")",
              "规则约束: 明确输出格式、字数、风格要求",
              "Few-shot 示例: 诗词模板包含完整示例和字数验证说明",
              "附件注入: 上传文件文本追加到 prompt 末尾作为参考资料",
              "风格指令: 通过 STYLE_MAP 动态注入风格修饰语",
          ], accent=PURPLE)

_add_card(slide, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.6),
          "量化与推理优化", [
              "GGUF 格式: 模型权重量化为 4-bit / 8-bit 整数",
              "KV Cache: 缓存已计算的 Key/Value 避免重复计算",
              "Flash Attention: 减少注意力计算的显存占用",
              "Continuous Batching: Ollama 自动批处理并发请求",
              "9B 参数 × 4-bit ≈ 5-6 GB 显存，消费级 GPU 可运行",
          ], accent=PINK)


# ══════════════════════════════════════════════════════════════
# Slide 9: SSE 流式通信
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "SSE 流式通信机制", "Server-Sent Events Streaming")

_add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.6),
              "为什么选择 SSE 而不是 WebSocket？",
              font_size=20, color=WHITE, bold=True)
_add_bullet_list(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(1.2), [
    "单向通信: LLM 生成是 请求->流式响应 模式, 不需要双向通道",
    "实现简单: 后端 StreamingResponse + 前端 fetch ReadableStream, 无需额外协议",
    "自动重连: 浏览器原生支持 SSE 断线重连 (虽然本项目用 fetch 手动实现)",
    "HTTP 兼容: 无需 WebSocket 升级, Nginx 反向代理配置更简单 (proxy_buffering off)",
], font_size=15)

# Data flow
flow_text = (
    "前端 fetch()                   后端 StreamingResponse              Ollama stream\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "POST /api/writing/stream  ──►  build_prompt()  ──►  POST /api/generate\n"
    "                                                     stream=true\n"
    "                               ◄── {response:\"你\"} ◄──  逐行 JSON\n"
    "◄── data: \"你\"\\n\\n  ◄──        ◄── {response:\"好\"} ◄──\n"
    "◄── data: \"好\"\\n\\n  ◄──        ◄── {done: true}   ◄──\n"
    "◄── data: [DONE]\\n\\n ◄──\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "ReadableStream.read()          async for token        aiter_lines()"
)
flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8))
flow_box.fill.solid()
flow_box.fill.fore_color.rgb = RGBColor(20, 20, 30)
flow_box.line.color.rgb = RGBColor(60, 60, 80)
_add_text_box(slide, Inches(1.1), Inches(4.35), Inches(11.2), Inches(2.5),
              flow_text, font_size=13, color=CYAN, font_name="Consolas")


# ══════════════════════════════════════════════════════════════
# Slide 10: 提示词工程
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "提示词工程", "Prompt Engineering in Practice")

templates = [
    ("通用生成",    "GENERATE_PROMPT",    "角色+要求+风格注入", ACCENT),
    ("古典诗词",    "POETRY_PROMPT",      "格律规则+示例+校验", GREEN),
    ("文本润色",    "POLISH_PROMPT",      "保持原意+提升质量",  ORANGE),
    ("文本翻译",    "TRANSLATE_PROMPT",   "准确+流畅+保持风格", PURPLE),
    ("文本摘要",    "SUMMARIZE_PROMPT",   "核心观点+1/3篇幅",  PINK),
    ("小红书",     "XIAOHONGSHU_PROMPT",  "emoji标题+口语化",  ACCENT),
    ("公众号",     "GONGZHONGHAO_PROMPT", "故事切入+金句加粗",  GREEN),
    ("头条",       "TOUTIAO_PROMPT",      "悬念标题+短段落",   ORANGE),
    ("短剧脚本",    "AI_DRAMA_PROMPT",    "场景/对白/反转",    PURPLE),
]

for i, (name, const, desc, color) in enumerate(templates):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.8) + row * Inches(1.6)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(1)

    _add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.3), Inches(0.35),
                  name, font_size=16, color=color, bold=True)
    _add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.3), Inches(0.3),
                  const, font_size=11, color=MEDIUM_GRAY, font_name="Consolas")
    _add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.3), Inches(0.35),
                  desc, font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 11: 特色功能
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)
_add_section_title(slide, "特色功能", "Key Features")

features = [
    ("诗词格律校验",
     "自动检测诗词请求 → 非流式生成 → 正则校验每句汉字数 → 不合格自动重试 (最多 3 次) → 逐字符 SSE 输出",
     ACCENT),
    ("智能文件解析",
     "PDF (PyPDF2 逐页提取) / Word (python-docx 段落提取) / 纯文本 (utf-8→gbk→gb2312→latin-1 多编码尝试)",
     GREEN),
    ("Word 文档导出",
     'Markdown 解析 -> python-docx 生成: 标题/列表/引用/加粗/斜体 -> 内存中生成 BytesIO -> 浏览器"另存为"对话框',
     ORANGE),
    ("历史记录搜索",
     "localStorage 持久化 (最多 50 条) → 关键词匹配内容和结果 → 按时间倒序排列 → 一键恢复历史结果",
     PURPLE),
]

for i, (title, desc, color) in enumerate(features):
    y = Inches(1.8) + i * Inches(1.35)
    # color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.08), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    _add_text_box(slide, Inches(1.3), y, Inches(11), Inches(0.35),
                  title, font_size=20, color=WHITE, bold=True)
    _add_text_box(slide, Inches(1.3), y + Inches(0.4), Inches(11), Inches(0.7),
                  desc, font_size=14, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 12: 总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide)

_add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
              "总结", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(2.4), Inches(1.7), Pt(3))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

summaries = [
    ("前后端分离", "FastAPI + React，职责清晰，独立部署"),
    ("本地化 LLM", "Ollama + Qwen 3.5，数据不出本机"),
    ("流式体验", "SSE 实时推送，毫秒级首字响应"),
    ("多场景覆盖", "9 种提示词模板，从学术到自媒体"),
    ("完整工具链", "上传解析 → AI 生成 → 导出 Word，端到端闭环"),
]

for i, (title, desc) in enumerate(summaries):
    y = Inches(2.9) + i * Inches(0.8)
    _add_text_box(slide, Inches(3.2), y, Inches(2.5), Inches(0.4),
                  title, font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(5.9), y, Inches(5), Inches(0.4),
                  desc, font_size=16, color=LIGHT_GRAY)

_add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
              "Thank You", font_size=24, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────
output = "/home/mystic/my_projects/my_first/docs/AI写作助手项目介绍.pptx"
prs.save(output)
print(f"PPT saved: {output}")
