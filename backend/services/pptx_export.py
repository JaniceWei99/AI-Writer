"""Convert AI-generated Markdown PPT outline to a .pptx file in memory.

Enhanced with decorative design elements for professional presentation style.

Expected input format (from PPT_PROMPT):
    ## 1. Slide Title
    - Bullet one
    - Bullet two
    > Speaker notes here
    ---
"""

import io
import re
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Template colour schemes ──────────────────────────────────────────

@dataclass
class PptTheme:
    """Color scheme and typography for a presentation template."""
    name: str
    bg: RGBColor
    title_color: RGBColor
    text_color: RGBColor
    accent: RGBColor
    accent_dim: RGBColor       # faded accent for large decorative shapes
    note_color: RGBColor
    bullet_color: RGBColor
    footer_bg: RGBColor        # footer/header bar color
    title_font: str = "Microsoft YaHei"
    body_font: str = "Microsoft YaHei"


THEMES: dict[str, PptTheme] = {
    "business": PptTheme(
        name="商务蓝",
        bg=RGBColor(0x0F, 0x17, 0x29),
        title_color=RGBColor(0x4F, 0xC3, 0xF7),
        text_color=RGBColor(0xDD, 0xDD, 0xDD),
        accent=RGBColor(0x29, 0x79, 0xFF),
        accent_dim=RGBColor(0x16, 0x25, 0x48),
        note_color=RGBColor(0x88, 0x88, 0x88),
        bullet_color=RGBColor(0x4F, 0xC3, 0xF7),
        footer_bg=RGBColor(0x12, 0x1C, 0x32),
    ),
    "minimal": PptTheme(
        name="简约白",
        bg=RGBColor(0xFA, 0xFA, 0xFA),
        title_color=RGBColor(0x1A, 0x1A, 0x1A),
        text_color=RGBColor(0x33, 0x33, 0x33),
        accent=RGBColor(0x33, 0x33, 0x33),
        accent_dim=RGBColor(0xE2, 0xE2, 0xE2),
        note_color=RGBColor(0x99, 0x99, 0x99),
        bullet_color=RGBColor(0x55, 0x55, 0x55),
        footer_bg=RGBColor(0xEE, 0xEE, 0xEE),
    ),
    "green": PptTheme(
        name="清新绿",
        bg=RGBColor(0x0D, 0x1F, 0x0D),
        title_color=RGBColor(0x66, 0xBB, 0x6A),
        text_color=RGBColor(0xD8, 0xE8, 0xD8),
        accent=RGBColor(0x2E, 0x7D, 0x32),
        accent_dim=RGBColor(0x14, 0x30, 0x16),
        note_color=RGBColor(0x80, 0xA0, 0x80),
        bullet_color=RGBColor(0x66, 0xBB, 0x6A),
        footer_bg=RGBColor(0x10, 0x28, 0x12),
    ),
    "warm": PptTheme(
        name="暖色调",
        bg=RGBColor(0x2E, 0x1A, 0x0E),
        title_color=RGBColor(0xFF, 0xB7, 0x4D),
        text_color=RGBColor(0xE8, 0xD8, 0xC8),
        accent=RGBColor(0xFF, 0x8F, 0x00),
        accent_dim=RGBColor(0x44, 0x2C, 0x14),
        note_color=RGBColor(0xA0, 0x90, 0x80),
        bullet_color=RGBColor(0xFF, 0xB7, 0x4D),
        footer_bg=RGBColor(0x38, 0x22, 0x10),
    ),
}

DEFAULT_TEMPLATE = "business"


# ── Slide data model ────────────────────────────────────────────────

@dataclass
class SlideData:
    title: str = ""
    bullets: list[str] = field(default_factory=list)
    notes: str = ""


# ── Parser ──────────────────────────────────────────────────────────

_HEADING_RE = re.compile(r"^#{1,3}\s+(?:\d+\.\s*)?(.+)")
_BULLET_RE = re.compile(r"^[-*+]\s+(.*)")
_ORDERED_RE = re.compile(r"^\d+\.\s+(.*)")
_NOTE_RE = re.compile(r"^>\s*(.*)")
_HR_RE = re.compile(r"^[-*_]{3,}$")


def _parse_slides(md_text: str) -> list[SlideData]:
    """Parse AI-generated PPT Markdown into a list of SlideData."""
    slides: list[SlideData] = []
    current: SlideData | None = None

    for raw_line in md_text.split("\n"):
        line = raw_line.strip()

        # Heading → new slide
        m = _HEADING_RE.match(line)
        if m:
            if current:
                slides.append(current)
            current = SlideData(title=m.group(1).strip())
            continue

        # Horizontal rule → slide separator (only if current exists)
        if _HR_RE.match(line):
            if current:
                slides.append(current)
                current = None
            continue

        if not current:
            # Content before first heading — start a slide
            if line:
                current = SlideData()
            else:
                continue

        # Speaker notes (blockquote)
        m = _NOTE_RE.match(line)
        if m:
            note_text = m.group(1).strip()
            # Remove prefix like "演讲备注："
            note_text = re.sub(r"^演讲备注[：:]\s*", "", note_text)
            if current.notes:
                current.notes += "\n" + note_text
            else:
                current.notes = note_text
            continue

        # Bullet point
        m = _BULLET_RE.match(line) or _ORDERED_RE.match(line)
        if m:
            current.bullets.append(m.group(1).strip())
            continue

        # Non-empty line → treat as bullet
        if line:
            current.bullets.append(line)

    if current:
        slides.append(current)

    return slides


# ── PPTX builder helpers ────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, shape_type, left, top, width, height, fill_color: RGBColor):
    """Add a shape with solid fill and no border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: int, color: RGBColor, bold: bool = False,
                 alignment=PP_ALIGN.LEFT, font_name: str = "Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _strip_inline_md(text: str) -> str:
    """Remove **bold** and *italic* markers for plain text display."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text


def _add_bullets_with_markers(slide, left, top, width, height,
                               bullets: list[str], theme: PptTheme):
    """Add bullet points with colored bullet markers for visual style."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for j, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()

        # Colored bullet marker
        marker = p.add_run()
        marker.text = "●  "
        marker.font.size = Pt(11)
        marker.font.color.rgb = theme.bullet_color
        marker.font.name = theme.body_font

        # Bullet text
        text_run = p.add_run()
        text_run.text = _strip_inline_md(bullet)
        text_run.font.size = Pt(18)
        text_run.font.color.rgb = theme.text_color
        text_run.font.name = theme.body_font

        p.space_after = Pt(14)
        p.line_spacing = Pt(30)

    return txBox


# ── Decorative elements ─────────────────────────────────────────────

def _add_cover_decorations(slide, theme: PptTheme, slide_w, slide_h):
    """Add decorative shapes to cover slide for visual impact."""
    # Large circle — top right (partially off-screen)
    _add_shape(slide, MSO_SHAPE.OVAL,
               slide_w - Inches(3.5), Inches(-2.0),
               Inches(6), Inches(6),
               theme.accent_dim)

    # Small circle — bottom left (partially off-screen)
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(-1.0), slide_h - Inches(2.8),
               Inches(3.5), Inches(3.5),
               theme.accent_dim)

    # Tiny accent dot — decorative detail
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(2.5), Inches(1.2),
               Inches(0.4), Inches(0.4),
               theme.accent_dim)

    # Bottom accent line — full width
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, slide_h - Inches(0.08), slide_w, Inches(0.08),
               theme.accent)


def _add_content_decorations(slide, theme: PptTheme, slide_w, slide_h):
    """Add decorative shapes to content slides for professional look."""
    # Left accent sidebar — full height
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, 0, Inches(0.12), slide_h,
               theme.accent)

    # Top accent line — full width
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, 0, slide_w, Pt(3),
               theme.accent)

    # Bottom footer line
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, slide_h - Inches(0.45), slide_w, Pt(1),
               theme.footer_bg)

    # Decorative circle — bottom right corner (partially off-screen)
    _add_shape(slide, MSO_SHAPE.OVAL,
               slide_w - Inches(2.5), slide_h - Inches(2.2),
               Inches(3.5), Inches(3.5),
               theme.accent_dim)


def _add_end_decorations(slide, theme: PptTheme, slide_w, slide_h):
    """Add decorative shapes to end/thank-you slide."""
    # Large circle — right side
    _add_shape(slide, MSO_SHAPE.OVAL,
               slide_w - Inches(5), Inches(1.0),
               Inches(6), Inches(6),
               theme.accent_dim)

    # Small circle — top left
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(-0.6), Inches(-0.6),
               Inches(2.2), Inches(2.2),
               theme.accent_dim)

    # Bottom accent line
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               0, slide_h - Inches(0.08), slide_w, Inches(0.08),
               theme.accent)


# ── Main conversion function ────────────────────────────────────────

def markdown_to_pptx(md_text: str, title: str = "",
                     template: str = DEFAULT_TEMPLATE,
                     images: dict[int, bytes] | None = None) -> io.BytesIO:
    """Convert AI-generated PPT markdown to a .pptx BytesIO buffer.

    *images* is an optional mapping of slide index → JPEG/PNG bytes.
    When provided, content slides use a left-text / right-image layout.
    """
    theme = THEMES.get(template, THEMES[DEFAULT_TEMPLATE])
    slides_data = _parse_slides(md_text)

    if not slides_data:
        # If parsing fails, create a single slide with the raw text
        slides_data = [SlideData(title=title or "演示文稿", bullets=[md_text[:500]])]

    if images is None:
        images = {}

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    MARGIN = Inches(0.9)
    CONTENT_W = SLIDE_W - 2 * MARGIN

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        _set_slide_bg(slide, theme.bg)

        slide_title = _strip_inline_md(sd.title)

        is_cover = (i == 0)
        is_end = (i == len(slides_data) - 1) and any(
            kw in slide_title for kw in ("谢谢", "致谢", "Q&A", "感谢", "Thank")
        )
        has_image = i in images

        if is_cover:
            # ── Decorations first (behind content) ──
            _add_cover_decorations(slide, theme, SLIDE_W, SLIDE_H)

            # Title — large, centered
            _add_textbox(
                slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1.5),
                slide_title, font_size=44, color=theme.title_color,
                bold=True, alignment=PP_ALIGN.CENTER, font_name=theme.title_font,
            )

            # Accent line under title
            line_w = Inches(3.33)
            line_left = (SLIDE_W - line_w) // 2
            _add_shape(slide, MSO_SHAPE.RECTANGLE,
                       line_left, Inches(3.5), line_w, Pt(4),
                       theme.accent)

            # Subtitle from first bullet
            if sd.bullets:
                subtitle = _strip_inline_md(sd.bullets[0])
                _add_textbox(
                    slide, MARGIN, Inches(3.9), CONTENT_W, Inches(1),
                    subtitle, font_size=20, color=theme.text_color,
                    alignment=PP_ALIGN.CENTER, font_name=theme.body_font,
                )

        elif is_end:
            # ── Decorations ──
            _add_end_decorations(slide, theme, SLIDE_W, SLIDE_H)

            # Title — centered
            _add_textbox(
                slide, MARGIN, Inches(2.5), CONTENT_W, Inches(1.5),
                slide_title, font_size=40, color=theme.title_color,
                bold=True, alignment=PP_ALIGN.CENTER, font_name=theme.title_font,
            )

            # Accent line
            line_w = Inches(2.5)
            line_left = (SLIDE_W - line_w) // 2
            _add_shape(slide, MSO_SHAPE.RECTANGLE,
                       line_left, Inches(3.8), line_w, Pt(3),
                       theme.accent)

            if sd.bullets:
                extra = "\n".join(_strip_inline_md(b) for b in sd.bullets)
                _add_textbox(
                    slide, MARGIN, Inches(4.2), CONTENT_W, Inches(1),
                    extra, font_size=18, color=theme.text_color,
                    alignment=PP_ALIGN.CENTER, font_name=theme.body_font,
                )

        else:
            # ── Content slide ──
            # Decorations first (behind content)
            _add_content_decorations(slide, theme, SLIDE_W, SLIDE_H)

            # Title
            _add_textbox(
                slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.9),
                slide_title, font_size=28, color=theme.title_color,
                bold=True, font_name=theme.title_font,
            )

            # Title underline
            _add_shape(slide, MSO_SHAPE.RECTANGLE,
                       MARGIN, Inches(1.3), Inches(2.2), Pt(3),
                       theme.accent)

            # Determine text area width based on image presence
            IMG_W = Inches(4.5)
            IMG_H = Inches(4.0)
            TEXT_W = (CONTENT_W - IMG_W - Inches(0.5)) if has_image else CONTENT_W

            # Bullets with colored markers
            if sd.bullets:
                _add_bullets_with_markers(
                    slide, MARGIN, Inches(1.6), TEXT_W, Inches(5),
                    sd.bullets, theme,
                )

            # Image (right side)
            if has_image:
                img_stream = io.BytesIO(images[i])
                img_left = SLIDE_W - MARGIN - IMG_W
                img_top = Inches(1.6)
                slide.shapes.add_picture(img_stream, img_left, img_top, IMG_W, IMG_H)

        # Speaker notes
        if sd.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = sd.notes

        # Page number (bottom right, skip cover)
        if not is_cover:
            _add_textbox(
                slide,
                SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.4),
                Inches(1.0), Inches(0.3),
                str(i + 1), font_size=10, color=theme.note_color,
                alignment=PP_ALIGN.RIGHT, font_name=theme.body_font,
            )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
