"""Extract text content from uploaded files (PDF, Word, TXT, PPTX, etc.)."""

import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".pptx", ".ppt",
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB


def extract_text(filename: str, file_bytes: bytes) -> str:
    """Read file bytes and return plain text content."""
    ext = Path(filename).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型: {ext}，支持: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    if len(file_bytes) > MAX_FILE_SIZE:
        raise ValueError("文件大小超过 10 MB 限制")

    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext in (".docx", ".doc"):
        return _extract_docx(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(file_bytes)
    else:
        return _extract_plain(file_bytes)


def _extract_pdf(data: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        parts.append(" | ".join(row_text))
        if parts:
            slides_text.append(f"--- 第{i}页 ---\n" + "\n".join(parts))
    return "\n\n".join(slides_text)


def _extract_plain(data: bytes) -> str:
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")
